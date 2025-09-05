from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
import numpy as np
from typing import Dict
import json
import logging
import io
import math

logger = logging.getLogger(__name__)

def create_pptx_from_json(slides_json: Dict, output_path: str) -> bool:
    """
    Create a PowerPoint presentation from JSON data.

    Args:
        slides_json: Dictionary containing slides data
        output_path: Path where to save the PowerPoint file

    Returns:
        bool: True if successful, False otherwise
    """
    try:
        logger.info(f"Creating PPTX from: {json.dumps(slides_json, indent=2)}")

        if not slides_json or "slides" not in slides_json:
            logger.error("Invalid slides data: missing 'slides' key")
            return False

        slides_data = slides_json.get("slides", [])
        if not slides_data:
            logger.error("No slides data provided")
            return False

        prs = Presentation()

        for i, slide_data in enumerate(slides_data):
            logger.info(f"Processing slide {i+1}: {json.dumps(slide_data, indent=2)}")

            # Use blank slide layout for better control over positioning
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Apply professional dark blue background by default
            background_color = slide_data.get("backgroundColor", "#1e3a8a")  # Default to professional dark blue
            if background_color and background_color != "#FFFFFF":
                try:
                    # Convert hex color to RGB
                    bg_color = background_color.lstrip('#')
                    r = int(bg_color[0:2], 16)
                    g = int(bg_color[2:4], 16)
                    b = int(bg_color[4:6], 16)

                    # Set slide background color
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(r, g, b)
                    logger.info(f"Applied background color {background_color} to slide {i+1}")
                except Exception as e:
                    logger.warning(f"Could not apply background color {background_color} to slide {i+1}: {str(e)}")
                    # Fallback to professional dark blue
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(30, 58, 138)  # Dark blue fallback

            # Get slide dimensions
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # Create title text box in top left corner
            title_left = int(slide_width * 0.05)  # 5% from left
            title_top = int(slide_height * 0.05)  # 5% from top
            title_width = int(slide_width * 0.9)  # 90% width
            title_height = int(slide_height * 0.15)  # 15% height

            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_paragraph = title_frame.paragraphs[0]

            # Set title with fallback
            title_text = slide_data.get("title", "").strip()
            if not title_text:
                title_text = f"Slide {i+1}"
                logger.warning(f"No title provided for slide {i+1}, using default: {title_text}")

            title_paragraph.text = title_text
            title_paragraph.font.size = Pt(32)  # Large title font
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White color for dark background

            # Create content text box below the title
            content_left = int(slide_width * 0.05)  # 5% from left
            content_top = int(slide_height * 0.25)  # 25% from top (below title)
            content_width = int(slide_width * 0.9)  # 90% width
            content_height = int(slide_height * 0.65)  # 65% height

            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # Build content from headline and content
            headline = slide_data.get("headline", "").strip()
            content_text = slide_data.get("content", "").strip()

            # Combine headline and content
            full_content = ""
            if headline:
                # Add headline as first paragraph
                headline_paragraph = content_frame.paragraphs[0]
                headline_paragraph.text = headline
                headline_paragraph.font.size = Pt(24)
                headline_paragraph.font.bold = True
                headline_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White for dark background
                full_content += f"{headline}\n\n"
            else:
                # Remove empty first paragraph if no headline
                if content_frame.paragraphs:
                    content_frame.paragraphs[0].text = ""

            if content_text:
                # Add content
                if headline:
                    # Add new paragraph for content
                    content_paragraph = content_frame.add_paragraph()
                    content_paragraph.text = content_text
                else:
                    # Use first paragraph for content
                    content_paragraph = content_frame.paragraphs[0]
                    content_paragraph.text = content_text

                content_paragraph.font.size = Pt(18)
                content_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White for dark background
                logger.info(f"Added content to slide {i+1}: {len(full_content)} characters")
            else:
                if not headline:
                    content_frame.paragraphs[0].text = "Slide content will appear here"
                    content_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White for dark background
                    logger.warning(f"No content provided for slide {i+1}")

            # Handle chart data if present
            chart_type = slide_data.get("chartType")
            chart_data = slide_data.get("chartData")

            if chart_type and chart_data:
                try:
                    chart_text = _process_chart_data(chart_type, chart_data)
                    if chart_text:
                        # Add chart info to content
                        chart_paragraph = content_frame.add_paragraph()
                        chart_paragraph.text = chart_text
                        chart_paragraph.font.size = Pt(16)
                        chart_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White for dark background
                        logger.info(f"Added chart data to slide {i+1}: {chart_type}")
                except Exception as e:
                    logger.error(f"Error processing chart for slide {i+1}: {str(e)}")
                    chart_paragraph = content_frame.add_paragraph()
                    chart_paragraph.text = f"\n\n[Chart: {chart_type} - Data visualization available]"
                    chart_paragraph.font.size = Pt(16)
                    chart_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White for dark background

        prs.save(output_path)
        logger.info(f"Successfully saved PowerPoint to: {output_path}")
        return True

    except Exception as e:
        logger.error(f"Error creating PPTX: {str(e)}")
        return False

def _process_chart_data(chart_type: str, chart_data) -> str:
    """
    Process chart data and return formatted text for the slide.

    Args:
        chart_type: Type of chart (bar, pie, line, etc.)
        chart_data: Chart data (dict or string)

    Returns:
        str: Formatted chart text for the slide
    """
    try:
        chart_text = f"\n\n📊 Chart Type: {chart_type.upper()}\n"

        if isinstance(chart_data, dict):
            labels = chart_data.get('labels', [])
            values = chart_data.get('values', [])
            chart_title = chart_data.get('title', f'{chart_type.title()} Chart')

            chart_text += f"Title: {chart_title}\n"
            chart_text += f"Data Points: {len(labels)}\n"

            # Show sample data
            if labels and values and len(labels) == len(values):
                chart_text += "\nSample Data:\n"
                for i, (label, value) in enumerate(zip(labels[:3], values[:3])):
                    chart_text += f"• {label}: {value}\n"
                if len(labels) > 3:
                    chart_text += f"• ... and {len(labels) - 3} more data points\n"
            else:
                chart_text += "Data format issue: labels and values arrays don't match\n"
        else:
            # Handle string data
            chart_text += f"Data: {str(chart_data)[:100]}..."
            if len(str(chart_data)) > 100:
                chart_text += "..."

        return chart_text

    except Exception as e:
        logger.error(f"Error processing chart data: {str(e)}")
        return f"\n\n[Chart processing error: {str(e)}]"
