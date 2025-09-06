import json
import io
import math
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
import numpy as np
from datetime import datetime

# --- Dark Mode Branding Constants ---
text_font = 'Arial'
key_font = 'Arial Bold'
heading_font = 'Arial Bold'

# --- Slide Dimension Constants ---
SLIDE_WIDTH = Inches(16)
SLIDE_HEIGHT = Inches(9)
SLIDE_MARGIN = Inches(0.5)
TITLE_HEIGHT = Inches(0.8)
CONTENT_TOP = Inches(1.2)
CONTENT_MAX_WIDTH = SLIDE_WIDTH - (2 * SLIDE_MARGIN)
CONTENT_MAX_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - SLIDE_MARGIN

SLIDE_BACKGROUND_COLOR = RGBColor(15, 22, 50)
DEFAULT_TEXT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEADER_BG_COLOR = RGBColor(0x44, 0x54, 0x6A)
TABLE_HEADER_FONT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
ROW_COLOR_DARK = RGBColor(0x2A, 0x39, 0x50)
BRAND_COLORS = ['#007ACC', '#09534F', '#4CAF50', '#FF9800', '#F44336', '#9C27B0']
HYPERLINK_COLOR = RGBColor(0x9B, 0xC1, 0xE4)

# --- Helper Functions ---
def set_title_style(title_shape, presentation_width):
    title_shape.left = Inches(0)
    title_shape.width = SLIDE_WIDTH
    title_shape.height = TITLE_HEIGHT
    title_shape.fill.background()
    line = title_shape.line
    line.color.rgb = TABLE_HEADER_BG_COLOR
    line.width = Pt(1)
    font = title_shape.text_frame.paragraphs[0].font
    font.name = heading_font
    font.size = Pt(28)
    font.color.rgb = DEFAULT_TEXT_COLOR
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf = title_shape.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)

def ensure_content_fits(left, top, width, height):
    """Ensure content stays within slide boundaries"""
    max_left = SLIDE_WIDTH - SLIDE_MARGIN - width
    max_top = SLIDE_HEIGHT - SLIDE_MARGIN - height
    
    left = max(SLIDE_MARGIN, min(left, max_left))
    top = max(SLIDE_MARGIN, min(top, max_top))
    
    return left, top, width, height

def calculate_chart_size():
    """Calculate optimal chart size"""
    return Inches(4), Inches(3)

def truncate_text_if_needed(text, max_length):
    """Truncate text to prevent overflow"""
    if len(text) <= max_length:
        return text
    return text[:max_length-3] + "..."

class GeneralPresentation:
    def __init__(self, data, search_phrase="Business Analysis"):
        if not data:
            raise ValueError("Input data is empty.")
        self.data = data
        self.search_phrase = search_phrase
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self.MAX_ROWS_PER_TABLE = 10
        
        # Set background for all layouts
        for layout in self.prs.slide_layouts:
            layout.background.fill.solid()
            layout.background.fill.fore_color.rgb = SLIDE_BACKGROUND_COLOR
    
    def _set_cell_style(self, cell, text, is_header=False, is_dark_row=False):
        """Style table cells with professional formatting"""
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        font = p.font
        font.name = text_font
        font.size = Pt(12)
        
        if is_header:
            font.bold = True
            font.color.rgb = TABLE_HEADER_FONT_COLOR
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG_COLOR
        else:
            font.color.rgb = DEFAULT_TEXT_COLOR
            if is_dark_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_COLOR_DARK
        
        cell.text_frame.margin_left = Inches(0.1)
        cell.text_frame.margin_right = Inches(0.1)
        cell.text_frame.margin_top = Inches(0.05)
        cell.text_frame.margin_bottom = Inches(0.05)

    def _create_data_chart(self, chart_data, chart_type="bar"):
        """Create various types of charts from data"""
        if not chart_data or 'labels' not in chart_data or 'values' not in chart_data:
            return None
            
        labels = chart_data['labels']
        values = chart_data['values']
        
        fig, ax = plt.subplots(figsize=(8, 6))
        fig.patch.set_facecolor('none')
        ax.set_facecolor('none')
        
        if chart_type == "pie":
            wedges, texts, autotexts = ax.pie(
                values, labels=labels, autopct='%.1f%%', startangle=90,
                colors=BRAND_COLORS[:len(values)], textprops={'color': 'white'}
            )
            plt.setp(autotexts, size=10, weight="bold", fontname=key_font)
            plt.setp(texts, size=12, fontname=text_font)
        
        elif chart_type == "bar":
            bars = ax.bar(labels, values, color=BRAND_COLORS[:len(values)])
            ax.set_ylabel('Values', color='white')
            ax.set_xlabel('Categories', color='white')
            ax.tick_params(colors='white')
            
            # Add value labels on bars
            for bar, value in zip(bars, values):
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height,
                       f'{value}', ha='center', va='bottom', color='white', fontweight='bold')
        
        elif chart_type == "line":
            ax.plot(labels, values, marker='o', linewidth=3, markersize=8, 
                   color=BRAND_COLORS[0], markerfacecolor=BRAND_COLORS[1])
            ax.set_ylabel('Values', color='white')
            ax.set_xlabel('Categories', color='white')
            ax.tick_params(colors='white')
            ax.grid(True, alpha=0.3, color='white')
        
        ax.spines['bottom'].set_color('white')
        ax.spines['top'].set_color('white')
        ax.spines['right'].set_color('white')
        ax.spines['left'].set_color('white')
        
        plt.tight_layout()
        chart_buffer = io.BytesIO()
        plt.savefig(chart_buffer, format='png', bbox_inches='tight', 
                   transparent=True, facecolor='none')
        plt.close(fig)
        chart_buffer.seek(0)
        return chart_buffer

    def _create_data_table(self, slide, table_data, title="Data Table"):
        """Create a professional data table"""
        if not table_data or 'headers' not in table_data or 'rows' not in table_data:
            return None
            
        headers = table_data['headers']
        rows = table_data['rows']
        
        # Calculate table dimensions
        num_cols = len(headers)
        num_rows = min(len(rows) + 1, self.MAX_ROWS_PER_TABLE + 1)  # +1 for header
        
        table_width = CONTENT_MAX_WIDTH * 0.8
        table_height = Inches(0.4) * num_rows
        
        table_left = SLIDE_MARGIN + (CONTENT_MAX_WIDTH - table_width) / 2
        table_top = CONTENT_TOP + Inches(0.5)
        
        table_left, table_top, table_width, table_height = ensure_content_fits(
            table_left, table_top, table_width, table_height
        )
        
        # Create table
        table = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, 
                                     table_width, table_height).table
        
        # Set headers
        for i, header in enumerate(headers):
            self._set_cell_style(table.cell(0, i), header, is_header=True)
        
        # Set data rows
        for row_idx, row_data in enumerate(rows[:self.MAX_ROWS_PER_TABLE]):
            is_dark = row_idx % 2 == 0
            for col_idx, cell_data in enumerate(row_data[:num_cols]):
                self._set_cell_style(table.cell(row_idx + 1, col_idx), 
                                   str(cell_data), is_dark_row=is_dark)
        
        return table

    def add_title_slide(self, title=None, subtitle=None):
        """Add a professional title slide"""
        if not title:
            title = f"Business Analysis: {self.search_phrase}"
        if not subtitle:
            subtitle = f"Comprehensive Analysis & Strategic Insights"
            
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        
        # Style title
        title_shape = slide.shapes.title
        title_shape.text_frame.paragraphs[0].font.name = heading_font
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = DEFAULT_TEXT_COLOR
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Style subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text_frame.paragraphs[0].font.name = text_font
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = DEFAULT_TEXT_COLOR

    def add_content_slide(self, slide_data):
        """Add a content slide with text, charts, and tables"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Set title
        title = slide_data.get('title', 'Content Slide')
        slide.shapes.title.text = title
        set_title_style(slide.shapes.title, self.prs.slide_width)
        
        # Add main content text
        content = slide_data.get('content', '')
        headline = slide_data.get('headline', '')
        
        if content or headline:
            text_left = SLIDE_MARGIN
            text_top = CONTENT_TOP
            text_width = CONTENT_MAX_WIDTH * 0.6  # Leave space for charts
            text_height = Inches(3)
            
            text_left, text_top, text_width, text_height = ensure_content_fits(
                text_left, text_top, text_width, text_height
            )
            
            txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Add headline if present
            if headline:
                p = tf.paragraphs[0]
                p.text = headline
                p.font.name = key_font
                p.font.size = Pt(20)
                p.font.color.rgb = DEFAULT_TEXT_COLOR
                p.font.bold = True
                
                # Add content as new paragraph
                if content:
                    p2 = tf.add_paragraph()
                    p2.text = f"\n{content}"
                    p2.font.name = text_font
                    p2.font.size = Pt(16)
                    p2.font.color.rgb = DEFAULT_TEXT_COLOR
            else:
                p = tf.paragraphs[0]
                p.text = content
                p.font.name = text_font
                p.font.size = Pt(16)
                p.font.color.rgb = DEFAULT_TEXT_COLOR
        
        # Add chart if data is available
        chart_data = slide_data.get('chartData')
        chart_type = slide_data.get('chartType', 'bar')
        
        if chart_data:
            chart_image = self._create_data_chart(chart_data, chart_type)
            if chart_image:
                chart_width, chart_height = calculate_chart_size()
                chart_left = SLIDE_WIDTH - chart_width - SLIDE_MARGIN
                chart_top = CONTENT_TOP
                
                chart_left, chart_top, chart_width, chart_height = ensure_content_fits(
                    chart_left, chart_top, chart_width, chart_height
                )
                
                slide.shapes.add_picture(chart_image, chart_left, chart_top, width=chart_width)
        
        # Add table if data is available
        table_data = slide_data.get('tableData')
        if table_data:
            # Position table below text content
            table_top = CONTENT_TOP + Inches(3.5)
            self._create_data_table(slide, table_data)

    def add_summary_slide(self):
        """Add a summary/conclusion slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = 'Key Takeaways & Next Steps'
        set_title_style(slide.shapes.title, self.prs.slide_width)
        
        # Get slides data to create summary
        slides = self.data.get('slides', [])
        
        text_left = SLIDE_MARGIN
        text_top = CONTENT_TOP
        text_width = CONTENT_MAX_WIDTH
        text_height = CONTENT_MAX_HEIGHT
        
        txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Summary content
        p = tf.paragraphs[0]
        p.text = "Summary of Key Findings:"
        p.font.name = key_font
        p.font.size = Pt(20)
        p.font.color.rgb = DEFAULT_TEXT_COLOR
        p.font.bold = True
        
        # Add key points from slides
        for i, slide_data in enumerate(slides[:4]):  # Max 4 key points
            p_bullet = tf.add_paragraph()
            title = slide_data.get('title', f'Point {i+1}')
            p_bullet.text = f"• {title}: Strategic importance for business growth"
            p_bullet.font.name = text_font
            p_bullet.font.size = Pt(16)
            p_bullet.font.color.rgb = DEFAULT_TEXT_COLOR
        
        # Next steps
        p_next = tf.add_paragraph()
        p_next.text = "\nRecommended Next Steps:"
        p_next.font.name = key_font
        p_next.font.size = Pt(18)
        p_next.font.color.rgb = DEFAULT_TEXT_COLOR
        p_next.font.bold = True
        
        next_steps = [
            "Develop detailed implementation roadmap",
            "Allocate necessary resources and budget",
            "Establish key performance indicators",
            "Begin pilot program execution"
        ]
        
        for step in next_steps:
            p_step = tf.add_paragraph()
            p_step.text = f"• {step}"
            p_step.font.name = text_font
            p_step.font.size = Pt(16)
            p_step.font.color.rgb = DEFAULT_TEXT_COLOR

def create_general_presentation(data, search_phrase="Business Analysis"):
    """Main function to create a general business presentation"""
    try:
        presentation = GeneralPresentation(data, search_phrase)
        
        # Add title slide
        presentation.add_title_slide()
        
        # Add content slides
        slides = data.get('slides', [])
        for slide_data in slides:
            presentation.add_content_slide(slide_data)
        
        # Add summary slide
        if len(slides) > 1:
            presentation.add_summary_slide()
        
        return presentation.prs
        
    except Exception as e:
        print(f"Error creating presentation: {str(e)}")
        return None
